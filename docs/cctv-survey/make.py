# -*- coding: utf-8 -*-
"""제니엘 버전 CCTV 설치 위치 조사 보고서 생성"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

NAVY  = RGBColor(0x0F, 0x2A, 0x4A)
BLUE  = RGBColor(0x1B, 0x6C, 0xA8)
SKY   = RGBColor(0xC9, 0xD8, 0xE8)
ACC   = RGBColor(0xE8, 0x72, 0x0C)
INK   = RGBColor(0x2F, 0x3A, 0x46)
MUTED = RGBColor(0x6B, 0x76, 0x82)
LIGHT = RGBColor(0xEE, 0xF2, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE  = RGBColor(0xD3, 0xDB, 0xE4)

FONT = '맑은 고딕'
SW, SH = 10.831597, 7.5
LM = 0.26
CW = SW - LM * 2

prs = Presentation('base.pptx')
LAY_COVER = prs.slide_layouts[10]    # 제목 및 내용 (de-branded)
LAY_BODY  = prs.slide_layouts[11]    # 1_제목 슬라이드 (de-branded, 제니엘 푸터)


# ---------- helpers ----------
def style(run, size, bold=False, color=INK, spacing=None):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ('ea', 'cs'):
        el = rPr.find('{%s}%s' % (A, tag))
        if el is None:
            el = etree.SubElement(rPr, '{%s}%s' % (A, tag))
        el.set('typeface', FONT)
    if spacing:
        rPr.set('spc', str(spacing))
    return run


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tb, tf


def para(tf, first=False, align=PP_ALIGN.LEFT, space_before=0, space_after=0, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line:
        p.line_spacing = line
    return p


def rect(slide, l, t, w, h, fill=None, line_col=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line_col is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_col; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sp.adjustments[0] = radius
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = ''
    return sp


def page_title(slide, main, sub=None):
    tb, tf = textbox(slide, LM, 0.16, CW, 0.42, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True)
    style(p.add_run(), 19, True, NAVY).text = main
    if sub:
        style(p.add_run(), 11, False, MUTED).text = '    ' + sub
    return tb


def bullet_lines(tf, items, size=11.5, color=INK, gap=5, first=True):
    for i, it in enumerate(items):
        p = para(tf, first=(first and i == 0), space_after=gap, line=1.25)
        style(p.add_run(), size, False, BLUE).text = '· '
        if isinstance(it, tuple):
            style(p.add_run(), size, True, INK).text = it[0]
            style(p.add_run(), size, False, color).text = it[1]
        else:
            style(p.add_run(), size, False, color).text = it


# =====================================================================
# 1. 표지
# =====================================================================
cover = prs.slides[0]
for sh in list(cover.shapes):
    sh._element.getparent().remove(sh._element)

rect(cover, 0, 0, SW, SH, fill=NAVY)
rect(cover, 0, SH - 1.55, SW, 1.55, fill=RGBColor(0x0A, 0x1F, 0x38))

tb, tf = textbox(cover, 0.85, 0.72, 5.0, 0.34, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, first=True)
style(p.add_run(), 17, True, WHITE).text = '제니엘'
style(p.add_run(), 10, False, RGBColor(0x7E, 0x9C, 0xBD), spacing=200).text = '   ZENIEL'

tb, tf = textbox(cover, 0.85, 2.55, 8.9, 0.3)
p = para(tf, first=True)
style(p.add_run(), 12, True, ACC, spacing=100).text = '현장 실사 보고서'

tb, tf = textbox(cover, 0.85, 3.00, 9.2, 1.0)
p = para(tf, first=True, line=1.15)
style(p.add_run(), 38, True, WHITE).text = '현장 CCTV 설치 장소(위치) 조사'

tb, tf = textbox(cover, 0.85, 4.15, 8.6, 0.6)
p = para(tf, first=True, line=1.3)
style(p.add_run(), 14, False, SKY).text = '포장재창고(1·2층)·도크장 CCTV 8대 설치 위치 선정 결과 및 설치 검토사항'

meta = [('대 상 현 장', '대전공장 포장재창고 1·2층 / 도크장'),
        ('조 사 일 자', '2026. 08. 18.'),
        ('작    성', '제니엘 (ZENIEL)')]
x = 0.85
for label, val in meta:
    tb, tf = textbox(cover, x, SH - 1.05, 3.1, 0.62)
    p = para(tf, first=True, space_after=4)
    style(p.add_run(), 9, False, RGBColor(0x7E, 0x9C, 0xBD), spacing=150).text = label
    p = para(tf, line=1.2)
    style(p.add_run(), 11.5, True, WHITE).text = val
    x += 3.15


# =====================================================================
# 2. 조사 개요
# =====================================================================
s = prs.slides.add_slide(LAY_BODY)
page_title(s, '조사 개요', '조사 목적 · 배경 · 범위 · 방법')

# KPI strip
kpis = [('총 설치 수량', '8', '대'), ('1층 포장재창고', '6', '대'),
        ('도크장(외부)', '1', '대'), ('2층 포장재창고', '1', '대')]
kw = (CW - 0.24 * 3) / 4
x = LM
for i, (label, num, unit) in enumerate(kpis):
    fill = NAVY if i == 0 else LIGHT
    rect(s, x, 0.82, kw, 0.98, fill=fill)
    tb, tf = textbox(s, x + 0.18, 0.95, kw - 0.36, 0.22)
    p = para(tf, first=True)
    style(p.add_run(), 9.5, False, SKY if i == 0 else MUTED).text = label
    tb, tf = textbox(s, x + 0.18, 1.20, kw - 0.36, 0.45)
    p = para(tf, first=True)
    style(p.add_run(), 26, True, WHITE if i == 0 else NAVY).text = num
    style(p.add_run(), 12, False, SKY if i == 0 else MUTED).text = ' ' + unit
    x += kw + 0.24

cards = [
    ('01', '조사 목적',
     ['창고·도크 구역 내 사각지대 및 지게차 동선 교차 구간에 대한 상시 관제 체계 확보',
      '안전사고·물류 사고(파손·분실) 발생 시 신속한 사실 확인 및 원인 규명',
      '현장 안전 수칙 준수 여부 모니터링을 통한 재발 방지 활동 근거 확보']),
    ('02', '조사 배경',
     ['포장재창고는 지게차·작업자 동선이 상시 교차하고, 랙 적재로 시야 사각지대가 발생',
      '도크장은 외부 차량 상하차가 이루어져 인·물 동시 위험 구간으로 관리 필요',
      '기존 관제 범위 밖 구역에 대한 CCTV 추가 설치 위치 선정 요청']),
    ('03', '조사 범위',
     ['포장재창고 1층 및 도크장(외부) : 7개소',
      '포장재창고 2층 : 1개소',
      '총 8개소 / 8대 (설치 지점·화각 방향·전원 및 통신 여건 포함)']),
    ('04', '조사 방법',
     ['현장 실사 후 배치 도면과 대조하여 설치 지점 표기',
      '지게차 운행 동선·반사경 위치·스피드도어 출입부 기준 사각지대 확인',
      '설치 지점별 현장 사진 촬영 및 전원·통신 인입 가능 여부 확인']),
]
cw = (CW - 0.26) / 2
ch = 2.14
for i, (no, title, lines) in enumerate(cards):
    cx = LM + (cw + 0.26) * (i % 2)
    cy = 2.00 + (ch + 0.22) * (i // 2)
    rect(s, cx, cy, cw, ch, fill=WHITE, line_col=LINE)
    rect(s, cx, cy, 0.62, 0.62, fill=NAVY)
    tb, tf = textbox(s, cx, cy, 0.62, 0.62, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    style(p.add_run(), 15, True, WHITE).text = no
    tb, tf = textbox(s, cx + 0.8, cy + 0.18, cw - 1.0, 0.3)
    p = para(tf, first=True)
    style(p.add_run(), 13.5, True, NAVY).text = title
    tb, tf = textbox(s, cx + 0.26, cy + 0.72, cw - 0.52, ch - 0.92)
    bullet_lines(tf, lines, size=10, gap=6)

tb, tf = textbox(s, LM, 6.62, CW, 0.24)
p = para(tf, first=True)
style(p.add_run(), 8.5, False, MUTED).text = (
    '※ 조사 목적·배경 항목은 현장 실사 결과를 바탕으로 제니엘이 작성한 초안이며, 발주처 확인 후 확정합니다.')


# =====================================================================
# 3. 설치 위치 총괄
# =====================================================================
s = prs.slides.add_slide(LAY_BODY)
page_title(s, '설치 위치 총괄', '총 8대 / 포장재창고 1층 6대 · 도크장(외부) 1대 · 포장재창고 2층 1대')

headers = ['No', '구분 (층·구역)', '설치 위치', '주요 감시 구역(안)', '비 고']
rows = [
    ('1', '포장재창고 1층', '외주사 휴게실 앞', '휴게실 전면 통로 및 작업자 출입 동선', '-'),
    ('2', '포장재창고 1층', '휴게실 외벽 상단', '창고 주통로 (남측) 지게차 운행 구간', '-'),
    ('3', '포장재창고 1층', '3Z-06 옆 기둥 근처', '3Z 랙 열 적재부 및 피킹 작업 구역', '-'),
    ('4', '포장재창고 1층', '측면 창가 위', '측면 통로 및 랙 사이 교차 동선', '-'),
    ('5', '포장재창고 1층', '스피드도어 위', '스피드도어 출입부 (지게차 진·출입)', '-'),
    ('6', '포장재창고 1층', '반사경 옆', '반사경 설치 교차로 (사각지대 구간)', '-'),
    ('7', '도크장 (외부)', '원료창고 위', '도크 접안부 및 상·하차 작업 구역', '모뎀 추가 연동'),
    ('8', '포장재창고 2층', '스피드도어 우측', '입식 지게차 운행구역', '-'),
]
colw = [0.55, 1.75, 2.15, 3.86, 2.0]
th, rh = 0.40, 0.415
tbl_shape = s.shapes.add_table(len(rows) + 1, 5, Inches(LM), Inches(1.02),
                               Inches(sum(colw)), Inches(th + rh * len(rows)))
tbl = tbl_shape.table
tbl.first_row = True
tbl.horz_banding = False
# 기본 테이블 스타일 제거(자체 서식 적용)
tblPr = tbl._tbl.find('{%s}tblPr' % A)
for st in tblPr.findall('{%s}tableStyleId' % A):
    tblPr.remove(st)

for i, w in enumerate(colw):
    tbl.columns[i].width = Inches(w)
tbl.rows[0].height = Inches(th)
for r in range(1, len(rows) + 1):
    tbl.rows[r].height = Inches(rh)

def fill_cell(cell, text, size, bold, color, bg, align=PP_ALIGN.LEFT):
    cell.fill.solid(); cell.fill.fore_color.rgb = bg
    cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    style(p.add_run(), size, bold, color).text = text

for c, h in enumerate(headers):
    fill_cell(tbl.cell(0, c), h, 11, True, WHITE, NAVY,
              PP_ALIGN.CENTER if c in (0, 4) else PP_ALIGN.LEFT)
for r, row in enumerate(rows, start=1):
    bg = WHITE if r % 2 else LIGHT
    for c, val in enumerate(row):
        is_dock = (row[4] != '-' and c == 4)
        fill_cell(tbl.cell(r, c), val, 10.5,
                  bold=(c == 0 or is_dock),
                  color=(ACC if is_dock else (NAVY if c == 0 else INK)),
                  bg=bg,
                  align=PP_ALIGN.CENTER if c in (0, 4) else PP_ALIGN.LEFT)

notes = [
    '외부 도크장 좌·우측 CCTV는 유선 인입이 어려워 모뎀을 추가하여 연동 실시합니다.',
    '주요 감시 구역(안)은 현장 실사 사진을 기준으로 제니엘이 작성한 초안이며, 발주처 확인 후 확정합니다.',
]
ny = 1.02 + th + rh * len(rows) + 0.22
rect(s, LM, ny, CW, 0.86, fill=LIGHT)
tb, tf = textbox(s, LM + 0.22, ny + 0.16, CW - 0.44, 0.56)
bullet_lines(tf, notes, size=10)

# 확인 필요사항
cy2 = ny + 1.04
rect(s, LM, cy2, CW, 0.72, fill=RGBColor(0xFD, 0xF3, 0xE8), line_col=RGBColor(0xF0, 0xC9, 0x9A))
tb, tf = textbox(s, LM + 0.22, cy2 + 0.13, CW - 0.44, 0.5)
p = para(tf, first=True, space_after=3)
style(p.add_run(), 10.5, True, ACC).text = '[확인 필요] '
style(p.add_run(), 10.5, True, INK).text = '도크장(외부) 설치 수량'
p = para(tf, line=1.2)
style(p.add_run(), 10, False, INK).text = (
    '총괄 수량은 도크장 외부 1대로 기재되어 있으나, 모뎀 연동 대상은 "좌·우측 2대"로 표기되어 있어 '
    '최종 수량(총 8대 / 9대) 확정이 필요합니다.')


# =====================================================================
# 6. 설치 시 검토사항  (도면 2장 뒤)
# =====================================================================
s = prs.slides.add_slide(LAY_BODY)
page_title(s, '설치 시 검토사항', '설치 확정 전 현장 여건 확인이 필요한 항목')

blocks = [
    ('전원 · 통신', NAVY,
     ['설치 지점별 전원 인입 위치 및 여유 용량 확인 (분전반 거리 / 배선 경로)',
      '도크장 외부는 유선 인입이 어려워 모뎀(무선) 추가 연동 방식 적용',
      '녹화 장비(NVR) 설치 위치 및 기존 관제 시스템 연동 방식 확정']),
    ('화각 · 사각지대', BLUE,
     ['랙 적재 높이 변동에 따른 시야 가림 여부 사전 검토 (만적 기준 화각 확인)',
      '반사경 설치 교차로는 양방향 동선이 함께 잡히도록 설치 각도 조정',
      '스피드도어·도크 접안부는 역광 발생 구간으로 노출 보정 기능 필요']),
    ('설치 환경', BLUE,
     ['스피드도어 상부는 개폐 진동이 발생하므로 브래킷 고정 방식 보강',
      '창고 내 분진·온습도 및 외부 지점의 방수·방진 등급(IP66 이상 권장) 확인',
      '조도 편차가 큰 구간은 야간·소등 시간대 촬영 품질 사전 확인']),
    ('설치 작업 안전', ACC,
     ['고소 작업 구간으로 작업허가서 발행 및 안전대·고소작업대 사용 필수',
      '지게차 운행 구역 내 작업 시 구역 통제 및 유도자 배치',
      '창고 운영 시간을 고려하여 비가동 시간대 작업 우선 협의']),
]
bw = (CW - 0.26) / 2
bh = 2.56
for i, (title, col, lines) in enumerate(blocks):
    bx = LM + (bw + 0.26) * (i % 2)
    by = 0.84 + (bh + 0.28) * (i // 2)
    rect(s, bx, by, bw, bh, fill=WHITE, line_col=LINE)
    rect(s, bx, by, bw, 0.5, fill=col)
    tb, tf = textbox(s, bx + 0.26, by, bw - 0.52, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True)
    style(p.add_run(), 13, True, WHITE).text = title
    tb, tf = textbox(s, bx + 0.26, by + 0.68, bw - 0.52, bh - 0.86)
    bullet_lines(tf, lines, size=10.5, gap=7)

tb, tf = textbox(s, LM, 6.62, CW, 0.24)
p = para(tf, first=True)
style(p.add_run(), 8.5, False, MUTED).text = (
    '※ 상기 항목은 현장 실사 기준 제니엘 검토(안)이며, 설치 업체 실측 및 발주처 협의를 거쳐 최종 확정합니다.')


# =====================================================================
# 7. 개인정보보호 준수사항
# =====================================================================
s = prs.slides.add_slide(LAY_BODY)
page_title(s, '개인정보보호 준수사항', '영상정보처리기기 설치·운영 시 사전 조치 사항')

items = [
    ('안내판 설치', '개인정보 보호법 제25조',
     ['촬영 구역 내 정보주체가 쉽게 인식할 수 있는 위치에 안내판 설치',
      '기재 사항 : 설치 목적 및 장소 / 촬영 범위 및 시간 / 관리책임자 성명·연락처']),
    ('근로자 대표 협의', '근로자참여 및 협력증진에 관한 법률 제20조',
     ['"사업장 내 근로자 감시 설비의 설치"는 노사협의회 협의 사항에 해당',
      '설치 목적이 근로자 감시가 아님을 명확히 하고, 사전 협의 및 설명 실시']),
    ('운영·관리 방침', '보관 및 접근 통제',
     ['영상정보처리기기 운영·관리 방침 수립 및 공개 (보관 기간 명시)',
      '열람 권한자 지정, 열람 이력 관리, 보관 기간 경과 시 자동 파기 설정']),
    ('도급 현장 운영', '제니엘 준수 사항',
     ['도급 인력의 영상 열람은 발주처 승인 절차를 거친 경우로 한정',
      '입사 교육 시 CCTV 설치 목적·범위 안내 및 확인 서명 관리']),
]
ry = 0.88
rh2 = 1.24
for i, (title, tag, lines) in enumerate(items):
    y = ry + (rh2 + 0.14) * i
    rect(s, LM, y, CW, rh2, fill=WHITE if i % 2 == 0 else LIGHT, line_col=LINE)
    rect(s, LM, y, 3.05, rh2, fill=NAVY if i % 2 == 0 else RGBColor(0x1B, 0x3D, 0x63))
    tb, tf = textbox(s, LM + 0.24, y + 0.22, 2.62, 0.82)
    p = para(tf, first=True, space_after=4)
    style(p.add_run(), 13, True, WHITE).text = title
    p = para(tf, line=1.2)
    style(p.add_run(), 8, False, SKY).text = tag
    tb, tf = textbox(s, LM + 3.34, y + 0.26, CW - 3.62, rh2 - 0.46)
    bullet_lines(tf, lines, size=10.5, gap=7)

tb, tf = textbox(s, LM, 6.62, CW, 0.24)
p = para(tf, first=True)
style(p.add_run(), 8.5, False, MUTED).text = (
    '※ 법령 조항은 참고용으로 정리한 것이며, 최종 적용 범위는 발주처 개인정보보호 담당 부서 검토 결과에 따릅니다.')


# =====================================================================
# 8. 추진 일정 및 협조 요청
# =====================================================================
s = prs.slides.add_slide(LAY_BODY)
page_title(s, '추진 일정 및 협조 요청', '위치 확정 이후 단계별 추진 계획(안)')

steps = [
    ('STEP 1', '설치 위치 확정', ['총괄표 기준 8개소 승인', '도크장 외부 수량 확정']),
    ('STEP 2', '자재 발주 · 입고', ['카메라·모뎀·NVR 사양 확정', '설치 업체 선정 및 실측']),
    ('STEP 3', '설치 · 배선 공사', ['전원/통신 배선 시공', '고소작업 허가 및 구역 통제']),
    ('STEP 4', '시운전 · 운영 개시', ['화각 조정 및 녹화 확인', '안내판 설치 · 관제 교육']),
]
sw_ = (CW - 0.24 * 3) / 4
for i, (st, title, lines) in enumerate(steps):
    x = LM + (sw_ + 0.24) * i
    rect(s, x, 0.90, sw_, 1.98, fill=WHITE, line_col=LINE)
    rect(s, x, 0.90, sw_, 0.36, fill=NAVY if i == 0 else BLUE)
    tb, tf = textbox(s, x, 0.90, sw_, 0.36, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    style(p.add_run(), 10, True, WHITE, spacing=100).text = st
    tb, tf = textbox(s, x + 0.2, 1.42, sw_ - 0.4, 0.32)
    p = para(tf, first=True)
    style(p.add_run(), 12.5, True, NAVY).text = title
    tb, tf = textbox(s, x + 0.2, 1.86, sw_ - 0.4, 0.9)
    bullet_lines(tf, lines, size=10, gap=6)
    tb, tf = textbox(s, x + 0.2, 2.52, sw_ - 0.4, 0.28)
    p = para(tf, first=True)
    style(p.add_run(), 9, False, MUTED).text = '일정 : (      .      ~      .      )'

# 협조 요청
rect(s, LM, 3.14, CW, 2.05, fill=WHITE, line_col=LINE)
rect(s, LM, 3.14, CW, 0.46, fill=NAVY)
tb, tf = textbox(s, LM + 0.26, 3.14, CW - 0.52, 0.46, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, first=True)
style(p.add_run(), 13, True, WHITE).text = '발주처 협조 요청 사항'
reqs = [
    ('전원 · 통신 인입 지원 : ', '설치 지점별 전원 확보 위치 및 사내망(또는 모뎀) 연동 조건 안내'),
    ('작업 여건 협조 : ', '고소작업 허가 발행, 지게차 운행 구역 통제, 비가동 시간대 작업 승인'),
    ('운영 기준 확정 : ', '영상 보관 기간, 열람 권한자 및 관리책임자 지정, 안내판 문안 확정'),
    ('설치 위치 최종 승인 : ', '총괄표 8개소 및 도크장 외부 수량(1대/2대) 확정'),
]
tb, tf = textbox(s, LM + 0.3, 3.76, CW - 0.6, 1.3)
for i, (h, body) in enumerate(reqs):
    p = para(tf, first=(i == 0), space_after=7, line=1.2)
    style(p.add_run(), 11, False, BLUE).text = '· '
    style(p.add_run(), 11, True, INK).text = h
    style(p.add_run(), 11, False, INK).text = body

# 마무리 배너
rect(s, LM, 5.42, CW, 1.02, fill=NAVY)
tb, tf = textbox(s, LM + 0.32, 5.58, CW - 0.64, 0.72)
p = para(tf, first=True, space_after=4)
style(p.add_run(), 12, True, WHITE).text = '제니엘은 설치 위치 확정 이후 현장 안전 관리와 운영 기준 준수를 함께 지원하겠습니다.'
p = para(tf, line=1.2)
style(p.add_run(), 10, False, SKY).text = (
    '작성 : 제니엘   |   본 자료는 현장 실사 결과를 정리한 보고용 자료이며, 발주처 확인 후 확정됩니다.')


# =====================================================================
# 슬라이드 순서 재배치 : 표지, 개요, 총괄, [도면1], [도면2], 검토, 개인정보, 일정
# =====================================================================
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)          # 0:표지 1:도면1 2:도면2 3:개요 4:총괄 5:검토 6:개인정보 7:일정
order = [0, 3, 4, 1, 2, 5, 6, 7]
for el in ids:
    sldIdLst.remove(el)
for i in order:
    sldIdLst.append(ids[i])

# 기존 도면 슬라이드 : 위치 라벨 줄바꿈(글자 겹침) 정리
LBL_MIN, LBL_MAX = 0.05, SW - 0.05
for idx in (3, 4):
    for sh in prs.slides[idx].shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text
        if '위치' not in txt or len(sh.text_frame.paragraphs) != 1:
            continue
        w_in = sh.width / 914400
        if w_in >= 1.7:
            continue
        center = sh.left / 914400 + w_in / 2
        new_l = min(max(center - 0.85, LBL_MIN), LBL_MAX - 1.7)
        sh.left, sh.width = Inches(new_l), Inches(1.7)
        if sh.shape_id in (137, 166):          # 사진 하단과 맞닿던 라벨 하향 조정
            sh.top = Inches(sh.top / 914400 + 0.10)

# 2층 도면 슬라이드 좌측 여백에 설치 정보 카드 추가
s2 = prs.slides[4]
rect(s2, 0.30, 0.92, 2.02, 1.52, fill=WHITE, line_col=LINE)
rect(s2, 0.30, 0.92, 2.02, 0.38, fill=NAVY)
tb, tf = textbox(s2, 0.30, 0.92, 2.02, 0.38, anchor=MSO_ANCHOR.MIDDLE)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
style(p.add_run(), 10.5, True, WHITE).text = '설치 No. 8'
tb, tf = textbox(s2, 0.48, 1.42, 1.68, 0.96)
for i, (k, v) in enumerate([('구분', '포장재창고 2층'), ('위치', '스피드도어 우측'),
                            ('감시', '입식 지게차 운행구역')]):
    p = para(tf, first=(i == 0), space_after=5, line=1.15)
    style(p.add_run(), 8.5, True, BLUE).text = k + '  '
    style(p.add_run(), 9, False, INK).text = v

# 기존 도면 슬라이드 제목 서식 통일
for idx in (3, 4):
    sl = prs.slides[idx]
    for sh in sl.shapes:
        if sh.has_text_frame and sh.name == '제목 2':
            sh.left, sh.top = Inches(LM), Inches(0.16)
            sh.width, sh.height = Inches(CW), Inches(0.42)
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.text = r.text.replace('□ ', '').replace('□', '')
                    style(r, 19, True, NAVY)
            break

prs.save('CCTV_설치위치_조사_제니엘.pptx')
print('saved')
